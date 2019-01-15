import unittest
from pptx import Presentation

class TestConvertPresentation(unittest.TestCase):

    def test_read_number_of_slides(self):
        prs = Presentation("testinputs/CK20V2.pptx")
        self.assertEqual(len(prs.slides), 100)


if __name__ == '__main__':
    unittest.main()
