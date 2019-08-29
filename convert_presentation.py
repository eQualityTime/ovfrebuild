from pptx import Presentation

def get_name(slide):
    for shape in slide.shapes:
      if shape.is_placeholder:
         if shape.placeholder_format.idx == 0:
             return shape.text
    return "" #Not an exception because this is the *name*, which isn't needed for OBZ


def is_first_inside_second(first,second):
    "Finds the centre the first shape and checks if it is in the bounds of the second"
    y=first['top']+(first['height']/2)-second['top']
    x=first['left']+(first['width']/2)-second['left']
    if 0 < y < second['height']:
        if 0 < x <second['width']:
            return True 
    return False


def hollow_shape(temp):
    hollow_shape={}
    hollow_shape['top']=temp.top
    hollow_shape['left']=temp.left
    hollow_shape['height']=temp.height
    hollow_shape['width']=temp.width
    return hollow_shape
    
class TreeNode: 
# A tree structure that stores the containers. 
    def __init__(self,slide_num,slide): 
        self.shape=hollow_shape(slide.shapes[slide_num])
        self.slide_num=slide_num
        #self top and left should be the same as shape.
        self.topleft=None
        self.topright=None
        self.bottomleft=None
        self.bottomright=None


    def add(self,slide_num,slide):
        newshape=hollow_shape(slide.shapes[slide_num])
        #print "entering add"
        number_of_additions=0
        if is_first_inside_second(newshape,self.shape):
            #print "Discarding for now but should fix later"
            return #fix this later, but we disgard things inside us for now.   
        else: 
                #print "New['top']: {},['left']: {}".format(newshape['top'],newshape['left'])
                #print "Current['top']: {},['left']: {}".format(self.shape['top'],self.shape['left'])
                if newshape['left']+newshape['width'] > self.shape['left']:
                    if newshape['top']+newshape['height']>self.shape['top']:
                        self.add_bottomright(slide_num, slide)
                        number_of_additions+=1
                    if newshape['top']<self.shape['top']:
                        self.add_topright(slide_num, slide)
                        number_of_additions+=1
                if newshape['left'] < self.shape['left']:
                    if newshape['top']+newshape['height']>self.shape['top']:
                        self.add_bottomleft(slide_num, slide)
                        number_of_additions+=1
                    if newshape['top']<self.shape['top']:
                        self.add_topleft(slide_num, slide)
                        number_of_additions+=1
        #print "Number of additions: {}".format(number_of_additions)
                    
    def add_bottomright(self,shape,slide):
        #print "br"
        temp=TreeNode(shape,slide)
        if self.bottomright==None:
            self.bottomright=temp
        else:
            self.bottomright.add(shape,slide)

    def add_bottomleft(self,shape,slide):
        #print "bl"
        temp=TreeNode(shape,slide)
        if self.bottomleft==None:
            self.bottomleft=temp
        else:
            self.bottomleft.add(shape,slide)
        
    def add_topright(self,shape,slide):
        #print "tr"
        temp=TreeNode(shape,slide)
        if self.topright==None:
            self.topright=temp
        else:
            self.topright.add(shape,slide)

    def add_topleft(self,shape,slide):
        #print "tl"
        temp=TreeNode(shape,slide)
        if self.topleft==None:
            self.topleft=temp
        else:
            self.topleft.add(shape,slide)
        

    def get_container_list(self):
# This needs a bit more work. 
        return_me=[]
        return_me.append(self.slide_num) 
        if self.topleft:
            return_me.extend(self.topleft.get_container_list())
        if self.topright:
            return_me.extend(self.topright.get_container_list())
        if self.bottomleft:
            return_me.extend(self.bottomleft.get_container_list())
        if self.bottomright:
            return_me.extend(self.bottomright.get_container_list())
        return return_me
                           


def get_containers(slide):
    root = TreeNode(0,slide)

    for i in range(len(slide.shapes))[1:]:
        #print "new shape"
        root.add(i,slide)

    containers=set(root.get_container_list() )
    return containers

    
