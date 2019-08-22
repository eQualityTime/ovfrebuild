#!/usr/bin/python
# -*- coding: utf-8 -*-
import unittest
from pptx import Presentation
import convert_presentation 
import time

class TestConvertPresentation(unittest.TestCase):

    def setUp(self):
        self._started_at = time.time()

    def tearDown(self):
        elapsed = time.time() - self._started_at
        print('{} ({}s)'.format(self.id(), round(elapsed, 2)))


    def test_read_number_of_slides(self):
        prs = Presentation("tests/testinputs/CK20V2.pptx")
        self.assertEqual(len(prs.slides), 104)

    def test_name(self): 
        slide = Presentation("tests/testinputs/CK20V2.pptx").slides[0]
        self.assertEqual(convert_presentation.get_name(slide), u"Top page")

    def test_name2(self): 
        slide = Presentation("tests/testinputs/CK20V2.pptx").slides[2]
        self.assertEqual(convert_presentation.get_name(slide), u"")

    def test_inside(self): 
        slide = Presentation("tests/testinputs/CK20V2.pptx").slides[3]
        shapeA=convert_presentation.hollow_shape(slide.shapes[1])
        shapeB=convert_presentation.hollow_shape(slide.shapes[0])
        self.assertEqual(True,convert_presentation.is_first_inside_second(shapeA,shapeB))

    def test_outside(self): 
        slide = Presentation("tests/testinputs/CK20V2.pptx").slides[4]
        shapeA=convert_presentation.hollow_shape(slide.shapes[1])
        shapeB=convert_presentation.hollow_shape(slide.shapes[0])
        self.assertEqual(False,convert_presentation.is_first_inside_second(shapeA,shapeB))
 
    def test_how_many_shapes_arent_inside_in_order(self):
        slide = Presentation("tests/testinputs/CK20V2.pptx").slides[2]
        containers=convert_presentation.get_containers(slide)
        self.assertEqual(8,len(containers))

    def test_tree_sorting(self):
        slide = Presentation("tests/testinputs/CK20V2.pptx").slides[3]
        containers=convert_presentation.get_containers(slide)
        self.assertEqual(1,len(containers))

    def test_tree_sorting2(self):
        slide = Presentation("tests/testinputs/CK20V2.pptx").slides[4]
        containers=convert_presentation.get_containers(slide)
        self.assertEqual(2,len(containers))

    def test_tree_sorting3(self):
        slide = Presentation("tests/testinputs/CK20V2.pptx").slides[5]
        containers=convert_presentation.get_containers(slide)
        self.assertEqual(4,len(containers))
    
    def test_get_all_containers(self):
        #mostly here for speed
        containers=[]
        for slide in Presentation("tests/testinputs/CK20V2.pptx").slides:
            containers=convert_presentation.get_containers(slide)
        self.assertEqual(23,len(containers))
        
        
#Todo - blank names 
        

if __name__ == '__main__':
    unittest.main()

pytest.main(__file__)
