from pptx import Presentation
import time

slide = Presentation("tests/testinputs/CK20V2.pptx").slides[1]
for i in range(1000):
    x=slide.shapes[0].top


#!/usr/bin/python
# -*- coding: utf-8 -*-
import unittest
from pptx import Presentation

class TestPPTX(unittest.TestCase):


    def setUp(self):
        self._started_at = time.time()

    def tearDown(self):
        elapsed = time.time() - self._started_at
        print('{} ({}s)'.format(self.id(), round(elapsed, 2)))

 
    def test_one(self):
        prs = Presentation("tests/testinputs/CK20V2.pptx")
        for i in range(100000):
            x = prs.slides[0].shapes[0].top
        
    def test_two(self):
        prs = Presentation("tests/testinputs/CK20V2.pptx")
        shape= prs.slides[0].shapes[0]
        for i in range(100000):
            x = shape.top

    def test_three(self):
        prs = Presentation("tests/testinputs/CK20V2.pptx")
        num= prs.slides[0].shapes[0].top
        for i in range(100000):
            x = num

if __name__ == '__main__':
    unittest.main()

